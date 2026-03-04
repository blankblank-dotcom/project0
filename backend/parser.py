"""
backend/parser.py — Multimodal Routing & Thumbnail Pipeline
==========================================================

Handles conversion of diverse document formats (PDF, PPTX, DOCX) 
into image sequences optimized for Qwen-VL (448x448).

Constraints:
- Memory Buffering: Uses io.BytesIO to avoid slow Disk I/O.
- Resolution: Native Qwen-VL size (448x448).
- Performance: Sequential processing for low RAM overhead.
"""

import io
import base64
import logging
from pathlib import Path
from typing import List, Union, Optional
from PIL import Image

log = logging.getLogger("parser")

# Target resolution for Qwen-VL native processing
QWEN_VL_SIZE = (448, 448)

def RouteToCanvas(file_path: Union[str, Path]) -> str:
    """Identify the document type and route to appropriate canvas handler.
    
    Returns:
        The extension type (lower case) if supported, else raises ValueError.
    """
    path = Path(file_path)
    suffix = path.suffix.lower()
    if suffix in [".pdf", ".pptx", ".docx"]:
        return suffix
    raise ValueError(f"Unsupported file format: {suffix}")

def convert_to_image_sequence(
    file_path: Union[str, Path], 
    target_size: tuple = QWEN_VL_SIZE
) -> List[io.BytesIO]:
    """Convert a document into a sequence of memory-buffered images.
    
    Args:
        file_path: Path to the document.
        target_size: Desired resolution (width, height).
        
    Returns:
        List of io.BytesIO objects containing PNG image data.
    """
    path = Path(file_path)
    ext = RouteToCanvas(path)
    
    if ext == ".pdf":
        return _pdf_to_images(path, target_size)
    elif ext == ".pptx":
        return _pptx_to_images(path, target_size)
    elif ext == ".docx":
        return _docx_to_images(path, target_size)
    
    return []

def get_preview_base64(file_path: Union[str, Path]) -> Optional[str]:
    """Returns the first page/slide as a base64 string for instant UI preview."""
    try:
        images = convert_to_image_sequence(file_path)
        if images:
            b64 = base64.b64encode(images[0].getvalue()).decode("utf-8")
            return f"data:image/png;base64,{b64}"
    except Exception as e:
        log.error(f"Failed to generate preview: {e}")
    return None

def _pdf_to_images(path: Path, target_size: tuple) -> List[io.BytesIO]:
    """Render PDF pages using PyMuPDF (fitz)."""
    import fitz
    doc = fitz.open(str(path))
    image_list = []
    
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        # Calculate matrix for target scale
        # Default PDF is 72 DPI. We want target_size.
        rect = page.rect
        zoom_x = target_size[0] / rect.width
        zoom_y = target_size[1] / rect.height
        matrix = fitz.Matrix(zoom_x, zoom_y)
        
        pix = page.get_pixmap(matrix=matrix, alpha=False)
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        buf.seek(0)
        image_list.append(buf)
        
    doc.close()
    return image_list

def _pptx_to_images(path: Path, target_size: tuple) -> List[io.BytesIO]:
    """Extract content from PPTX. 
    Note: Real rendering requires MS Office. 
    This falls back to a 'Digital Canvas' (Text + High-res Shapes).
    """
    from pptx import Presentation
    from PIL import ImageDraw, ImageFont
    
    prs = Presentation(str(path))
    image_list = []
    
    for slide in prs.slides:
        # Create a blank white canvas
        img = Image.new("RGB", target_size, color="white")
        draw = ImageDraw.Draw(img)
        
        # Simple heuristic: extract titles and main text
        y_offset = 20
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                # Draw text on canvas (simplified)
                text = shape.text.strip()[:100] # Limit per shape
                draw.text((20, y_offset), text, fill="black")
                y_offset += 30
                if y_offset > target_size[1] - 40:
                    break
        
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        buf.seek(0)
        image_list.append(buf)
        
    return image_list

def _docx_to_images(path: Path, target_size: tuple) -> List[io.BytesIO]:
    """Convert DOCX to images via 'Digital Canvas' logic."""
    from docx import Document
    from PIL import ImageDraw
    
    doc = Document(str(path))
    full_text = [p.text for p in doc.paragraphs if p.text.strip()]
    
    # Paginate text into canvases
    chars_per_page = 500
    pages_text = []
    current_page = ""
    
    for para in full_text:
        if len(current_page) + len(para) > chars_per_page:
            pages_text.append(current_page)
            current_page = para
        else:
            current_page += "\n" + para
    if current_page:
        pages_text.append(current_page)
        
    image_list = []
    for text_block in pages_text:
        img = Image.new("RGB", target_size, color="white")
        draw = ImageDraw.Draw(img)
        
        # Wrap and draw text
        y_offset = 20
        lines = text_block.split("\n")
        for line in lines:
            draw.text((20, y_offset), line[:50], fill="black")
            y_offset += 20
            if y_offset > target_size[1] - 30:
                break
                
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        buf.seek(0)
        image_list.append(buf)
        
    return image_list
